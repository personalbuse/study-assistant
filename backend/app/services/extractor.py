import hashlib
from pathlib import Path


def calculate_hash(filepath: str) -> str:
    sha256 = hashlib.sha256()
    with open(filepath, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            sha256.update(chunk)
    return sha256.hexdigest()


def extract_text(filepath: str) -> tuple[str, int]:
    ext = Path(filepath).suffix.lower()

    if ext == ".pdf":
        return extract_from_pdf(filepath)
    elif ext == ".pptx":
        return extract_from_pptx(filepath)
    elif ext == ".docx":
        return extract_from_docx(filepath)
    elif ext == ".md":
        return extract_from_txt(filepath)
    elif ext == ".txt":
        return extract_from_txt(filepath)
    else:
        raise ValueError(f"Formato no soportado: {ext}")


def extract_from_pdf(filepath: str) -> tuple[str, int]:
    import pdfplumber

    text_parts = []
    page_count = 0

    with pdfplumber.open(filepath) as pdf:
        page_count = len(pdf.pages)
        for page_num, page in enumerate(pdf.pages, 1):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(f"--- Página {page_num} ---\n{page_text}")

    return "\n\n".join(text_parts), page_count


def extract_from_pptx(filepath: str) -> tuple[str, int]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    text_parts = []
    prs = Presentation(filepath)
    slide_count = len(prs.slides)

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            _extract_pptx_shape_text(shape, slide_texts)
        if slide_texts:
            text_parts.append(
                f"--- Diapositiva {slide_num} ---\n" + "\n".join(slide_texts)
            )

    return "\n\n".join(text_parts), slide_count


def _extract_pptx_shape_text(shape, texts):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import pytesseract
    from PIL import Image
    import io

    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t:
                texts.append(t)
    if shape.has_table:
        for row in shape.table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                texts.append(" | ".join(cells))
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            image = Image.open(io.BytesIO(shape.image.blob))
            ocr_text = pytesseract.image_to_string(image, lang="spa").strip()
            if ocr_text:
                texts.append(ocr_text)
        except Exception:
            pass
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            _extract_pptx_shape_text(s, texts)


def extract_from_docx(filepath: str) -> tuple[str, int]:
    from docx import Document

    doc = Document(filepath)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text for cell in row.cells]
            paragraphs.append(" | ".join(row_text))

    return "\n".join(paragraphs), 1


def extract_from_txt(filepath: str) -> tuple[str, int]:
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    return text, 1


def extract_with_ocr(filepath: str) -> str:
    import pytesseract
    from pdf2image import convert_from_path

    images = convert_from_path(filepath)
    text_parts = []

    for page_num, image in enumerate(images, 1):
        page_text = pytesseract.image_to_string(image, lang="spa")
        text_parts.append(f"--- Página {page_num} (OCR) ---\n{page_text}")

    return "\n\n".join(text_parts)
